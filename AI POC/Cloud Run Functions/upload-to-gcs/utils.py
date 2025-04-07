import vertexai
from vertexai.generative_models import GenerativeModel, Part, SafetySetting
from constants import project_id,location,model_name
from retry import retry
import json
import logging
import os
import io

from google.auth import default
from google.cloud import storage, tasks_v2
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib.utils import simpleSplit
from textwrap import wrap

from constants import (
    gcs_bucket_name, 
    download_folder, 
    destination_folder, 
    metadata_folder, 
    metadata_file,
    project_id,
    location
)

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

### Gen AI based functions

def init_vertexai(project, location):
    """
    Initialize the Vertex AI with the specified project and location.

    Returns:
        GenerativeModel: The initialized generative model.
    """
    vertexai.init(project=project, location=location)
    return GenerativeModel(model_name)

def configure_generation():
    """
    Configure the generation parameters for the generative model.

    Returns:
        dict: A dictionary containing generation configuration.
    """
    return {
        "max_output_tokens": 8192,
        "temperature": 0,
        "top_p": 1,
    }

def configure_safety_settings():
    """
    Configure the safety settings for the generative model.

    Returns:
        list: A list of SafetySetting objects.
    """
    return [
        SafetySetting(
            category=SafetySetting.HarmCategory.HARM_CATEGORY_HATE_SPEECH,
            threshold=SafetySetting.HarmBlockThreshold.OFF
        ),
        SafetySetting(
            category=SafetySetting.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT,
            threshold=SafetySetting.HarmBlockThreshold.OFF
        ),
        SafetySetting(
            category=SafetySetting.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT,
            threshold=SafetySetting.HarmBlockThreshold.OFF
        ),
        SafetySetting(
            category=SafetySetting.HarmCategory.HARM_CATEGORY_HARASSMENT,
            threshold=SafetySetting.HarmBlockThreshold.OFF
        ),
    ]

@retry(tries=5, delay=2, backoff=2)
def generate_summary(content):
    """
    Generates a summary from the given content using the Gemini model.

    Args:
        content (str): The content of the article to summarize.

    Returns:
        str: The generated summary.
    """

    model = init_vertexai(project_id, location)
    generation_config = configure_generation()
    safety_settings = configure_safety_settings()

    prompt = f"""
    Please provide a comprehensive summary of the following article, ensuring that no key points are omitted.  Maintain the original writing style and tone as much as possible.

    Article:
    {content}
    """

    response = model.generate_content(
        prompt,
        generation_config=generation_config,
        safety_settings=safety_settings,
    )

    summary = response.text
    
    return summary

@retry(tries=5, delay=2, backoff=2)
def generate_title(content):
    """
    Generates a title from the given content using the Gemini model.

    Args:
        content (str): The content of the article to summarize.

    Returns:
        str: The generated summary.
    """

    model = init_vertexai(project_id, location)
    generation_config = configure_generation()
    safety_settings = configure_safety_settings()

    prompt = f"""
    Please provide a concise and descriptive title for the following article, capturing all key information.

    Article:
    {content}
    """

    response = model.generate_content(
        prompt,
        generation_config=generation_config,
        safety_settings=safety_settings,
    )

    title = response.text
    
    return title

### Auth based functions

def authenticate_drive():
    """Authenticates and returns a Google Drive service client."""
    credentials, project = default(scopes=["https://www.googleapis.com/auth/drive"])
    drive_service = build("drive", "v3", credentials=credentials)
    return drive_service

def authenticate_sheets():
    """Authenticates and returns a Google Sheets service client."""
    credentials, project = default(scopes=["https://www.googleapis.com/auth/spreadsheets"])
    sheets_service = build("sheets", "v4", credentials=credentials)
    return sheets_service

def authenticate_storage():
    """Authenticates and returns a Google Cloud Storage client."""
    credentials, project = default()
    storage_client = storage.Client(credentials=credentials, project=project)
    return storage_client



def upload_file_to_gcs(storage_client, bucket_name, source_file_path, file_name):
    """Uploads a file to Google Cloud Storage."""
    try:
        bucket = storage_client.get_bucket(bucket_name)
        destination_blob_name = f"{destination_folder}/{file_name}"
        blob = bucket.blob(destination_blob_name)
        blob.upload_from_filename(source_file_path)
        logger.info(f"Uploaded {source_file_path} to gs://{bucket_name}/{destination_blob_name}.")
        return source_file_path
    except Exception as e:
        logger.error(f"Error uploading the file to GCS bucket: {e}")
        return None


def extract_content_from_pdf(file_path,file_name):
    """
    Extracts text content from a PDF file in GCS.
    """
    try:
        with open(file_path, "rb") as file:
            reader = PdfReader(file)
            content = ""
            for page in reader.pages: 
                content += page.extract_text()
            return content  
    except Exception as e:
        logger.error(f"Error extracting text from PDF {file_name}: {e}")
        return None


def convert_docx_to_pdf(docx_path, pdf_path):
    """Convert .docx to .pdf with minimal margin, smaller font, and text wrapping."""
    try:
        doc = Document(docx_path)
        pdf_canvas = canvas.Canvas(pdf_path, pagesize=letter)
        
        x_margin = 30
        y_position = 780
        line_height = 15
        max_width = 500
        pdf_canvas.setFont("Helvetica", 10)

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                wrapped_lines = wrap(text, width=80)
                for line in wrapped_lines:
                    if y_position < 50:
                        pdf_canvas.showPage()
                        pdf_canvas.setFont("Helvetica", 10)
                        y_position = 780
                    
                    pdf_canvas.drawString(x_margin, y_position, line)
                    y_position -= line_height

        pdf_canvas.save()
        os.remove(docx_path)
    except Exception as e:
        logger.error(f"Exception : {e}")
        logger.error(f"Error File {docx_path}")

def convert_ppt_to_pdf(pptx_path,pdf_path):
    try:
        prs = Presentation(pptx_path)
        pdf_canvas = canvas.Canvas(pdf_path, pagesize=landscape(letter))

        for slide in prs.slides:
            y_position = 500
            pdf_canvas.setFont("Helvetica-Bold", 14)
            
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        wrapped_text = simpleSplit(text, "Helvetica-Bold", 14, 700)
                        for line in wrapped_text:
                            pdf_canvas.drawString(50, y_position, line)
                            y_position -= 25
                            if y_position < 50:
                                pdf_canvas.showPage()
                                pdf_canvas.setFont("Helvetica-Bold", 14)
                                y_position = 500

            pdf_canvas.showPage()

        pdf_canvas.save()
        os.remove(pptx_path)
        return pdf_path
    except Exception as e:
        logger.error(f"Exception : {e}")
        logger.error(f"Error File {pptx_path}")



def download_file_from_drive(drive_service, file_id, file_name, download_folder,mime):
    try:
        if mime == "application/vnd.google-apps.document":
            request = drive_service.files().export_media(fileId=file_id, mimeType="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        else:
            request = drive_service.files().get_media(fileId=file_id)
        file_path = os.path.join(download_folder,file_name)
        with io.BytesIO() as file_stream:
            downloader = MediaIoBaseDownload(file_stream, request)
        
            done = False
            while not done:
                _, done = downloader.next_chunk()
            with open(file_path, "wb") as f:
                f.write(file_stream.getvalue())
                
        return file_path
        
    except Exception as e:
        logger.error(f"Error with the drive file file Name: {file_name} {e}")
        return None

### Metadata functions

def create_metadata(id,title,mime,content,file_name) -> dict:
    """Creates a metadata dictionary with ID, uri, title, link, categories and summary"""

    uri = f'gs://{gcs_bucket_name}/{destination_folder}/{file_name}'
    link = f'https://storage.cloud.google.com/{gcs_bucket_name}/{destination_folder}/{file_name}'
    if len(content) > 1000 and len(content)<8000:
        content = generate_summary(content)
    elif len(content) > 8000:
        logger.info(f"{file_name} exceeds size for summary generation ")



    metadata = {
        'id': id,
        'content': {
            'mimeType': mime,
            
            'uri': uri
        },
        'structData': {
            'title': title,
            'link' : link,
            'content': content
        }
    }
    return metadata

def append_metadata_to_gcs(storage_client, metadata):
    """
    Appends metadata to the metadata.jsonl file in GCS.
    """
    bucket = storage_client.get_bucket(gcs_bucket_name)
    metadata_file_name = f"{metadata_folder}/{metadata_file}"
    blob = bucket.blob(metadata_file_name)

    try:
        if blob.exists():
            content = blob.download_as_text()
            metadata_lines = content.splitlines()
            updated_metadata_lines = []
            title = metadata['structData']['title']

            title_found= False

            for line in metadata_lines:
                existing_metadata = json.loads(line)
                if existing_metadata['structData']['title'] == title:
                    updated_metadata_lines.append(json.dumps(metadata))
                    title_found=True
                else:
                    updated_metadata_lines.append(line)

            if not title_found:
                updated_metadata_lines.append(json.dumps(metadata))
        else:
            updated_metadata_lines = [json.dumps(metadata)]

        blob.upload_from_string("\n".join(updated_metadata_lines))
        logger.info(f"Appended metadata to {metadata_file_name} in GCS.")

    except Exception as e:
        logger.error(f"Error appending metadata to {metadata_file_name}: {e}")
